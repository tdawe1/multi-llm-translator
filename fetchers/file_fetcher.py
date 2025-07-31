import docx
import pptx
import openpyxl


def get_text_from_file(filepath: str) -> str:
    """
    Reads and returns text content from a file, supporting .txt, .docx, .pptx, and .xlsx.
    """
    full_text_list = []
    try:
        if filepath.endswith(".docx"):
            doc = docx.Document(filepath)
            for para in doc.paragraphs:
                full_text_list.append(para.text)

        elif filepath.endswith(".pptx"):
            prs = pptx.Presentation(filepath)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        full_text_list.append(shape.text)

        elif filepath.endswith(".xlsx"):
            wb = openpyxl.load_workbook(filepath)
            for sheet in wb.worksheets:
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value and isinstance(cell.value, str):
                            full_text_list.append(cell.value)

        elif filepath.endswith(".txt"):
            with open(filepath, "r", encoding="utf-8") as f:
                return f.read()

        else:
            return "Error: Unsupported file type."

        # Use a consistent separator for regeneration
        return "\n\n".join(full_text_list)

    except FileNotFoundError:
        return f"Error: File not found at {filepath}"
    except Exception as e:
        return f"Error: Could not read file. {e}"
