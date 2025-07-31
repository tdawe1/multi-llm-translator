from pptx import Presentation


def create_pptx_from_text(
    original_filepath: str, translated_text: str, output_filepath: str
):
    """
    Creates a new .pptx file by replacing text in the original with translated text.
    """
    try:
        prs = Presentation(original_filepath)
        translated_paragraphs = translated_text.split("\n\n")
        para_index = 0

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if para_index < len(translated_paragraphs):
                        shape.text = translated_paragraphs[para_index]
                        para_index += 1

        prs.save(output_filepath)
        return f"Successfully created {output_filepath}"
    except Exception as e:
        return f"Error: Failed to regenerate .pptx file. {e}"
